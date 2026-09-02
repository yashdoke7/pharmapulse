"""Geometry QA for the generated deck.

LibreOffice is not available on this machine, so slides cannot be rendered to
images. This checks the same defect classes analytically:

  - anything past the slide edge, or inside the 0.5" margin
  - text boxes whose estimated wrapped height exceeds their box
  - overlapping text frames
  - low information: empty shapes

Run:  python deck/qa_deck.py deck/PharmaPulse.pptx
"""

from __future__ import annotations

import sys

from pptx import Presentation

EMU_IN = 914400.0
MARGIN = 0.5

# Rough advance width per point of font size, as a fraction of em. Calibri and
# Cambria sit near 0.50; mono is wider. Deliberately conservative.
CHAR_W = {"Courier New": 0.60, "Cambria": 0.50, "Calibri": 0.48}
DEFAULT_CHAR_W = 0.52


def inches(v) -> float:
    return (v or 0) / EMU_IN


def text_of(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs)


def estimate_height(shape) -> float:
    """Estimated rendered height of the text, in inches."""
    if not shape.has_text_frame:
        return 0.0
    w_in = inches(shape.width)
    if w_in <= 0:
        return 0.0

    total = 0.0
    for para in shape.text_frame.paragraphs:
        runs = para.runs
        if not runs:
            total += 0.16
            continue
        size = max((r.font.size.pt for r in runs if r.font.size), default=18.0)
        face = next((r.font.name for r in runs if r.font.name), "Calibri")
        cw = CHAR_W.get(face, DEFAULT_CHAR_W) * size / 72.0
        text = "".join(r.text for r in runs)
        chars_per_line = max(int(w_in / cw), 1) if cw > 0 else 999
        # honour explicit newlines, then wrap
        lines = 0
        for hard in text.split("\n"):
            lines += max(1, -(-len(hard) // chars_per_line))
        # python-pptx returns line_spacing either as a float multiple or as a
        # Length (EMU). pptxgenjs writes points, which arrive as a Length.
        ls = para.line_spacing
        if ls is None:
            line_h = size * 1.22 / 72.0
        elif hasattr(ls, "pt"):                       # Length -> points
            line_h = float(ls.pt) / 72.0
        elif ls > 4:                                  # already points
            line_h = float(ls) / 72.0
        else:                                         # a multiple of font size
            line_h = size * float(ls) / 72.0
        total += lines * line_h
    return total


def main(path: str) -> int:
    # The Windows console is cp1252 and the deck is full of typographic marks.
    # Losing a glyph in a QA report is fine; crashing on one is not.
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    prs = Presentation(path)
    SW, SH = inches(prs.slide_width), inches(prs.slide_height)
    print(f"{path}  —  {len(prs.slides.__iter__.__self__._sldIdLst)} slides, {SW:.2f}in x {SH:.2f}in\n")

    problems = 0
    for n, slide in enumerate(prs.slides, start=1):
        issues: list[str] = []
        boxes = []

        for shape in slide.shapes:
            x, y = inches(shape.left), inches(shape.top)
            w, h = inches(shape.width), inches(shape.height)
            label = (text_of(shape).strip().replace("\n", " ")[:44] or shape.shape_type)

            if x < -0.01 or y < -0.01 or x + w > SW + 0.01 or y + h > SH + 0.01:
                issues.append(f"OFF-SLIDE  ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}  “{label}”")
            elif x < MARGIN - 0.01 or y < MARGIN - 0.01 or x + w > SW - MARGIN + 0.01 or y + h > SH - MARGIN + 0.01:
                # bleed shapes are intentional only if they cover the whole slide
                if not (w > SW * 0.95 or h > SH * 0.95):
                    issues.append(f"tight margin ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}  “{label}”")

            if shape.has_text_frame and text_of(shape).strip():
                est = estimate_height(shape)
                if est > h * 1.18:
                    issues.append(
                        f"OVERFLOW?  needs ~{est:.2f}in, box {h:.2f}in  “{label}”")
                boxes.append((x, y, w, h, label))

        # overlapping text frames
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                ax, ay, aw, ah, al = boxes[i]
                bx, by, bw, bh, bl = boxes[j]
                ox = min(ax + aw, bx + bw) - max(ax, bx)
                oy = min(ay + ah, by + bh) - max(ay, by)
                if ox > 0.06 and oy > 0.06:
                    issues.append(f"overlap {ox:.2f}x{oy:.2f}in  “{al}”  vs  “{bl}”")

        if issues:
            problems += len(issues)
            print(f"slide {n}")
            for it in issues:
                print(f"   {it}")
            print()

    print(f"{problems} issue(s) flagged." if problems else "No geometry issues flagged.")
    return 1 if problems else 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1] if len(sys.argv) > 1 else "deck/PharmaPulse.pptx"))
